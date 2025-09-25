"""
BIP39 Seed Phrase Scanner
Scans files and directories for potential BIP39 seed phrases
"""

# First try to import all required dependencies
import sys
import os
import re
import json
import zipfile
import sqlite3
import logging
import string
import platform
import mimetypes
import chardet
import hashlib
from datetime import datetime, timedelta
from typing import Set, List, Dict, Optional, Tuple
from pathlib import Path
import magic  # python-magic for better mime type detection
import textract  # Added import for textract
import openpyxl  # Added import for openpyxl
from pptx import Presentation  # Added import for pptx
import time
from tqdm import tqdm

# Try to import PDF handling library
PDF_SUPPORT = False
try:
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

# Try to import DOCX handling library
DOCX_SUPPORT = False
try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

# Constants
ENCODINGS = ['utf-8', 'utf-16', 'utf-32', 'ascii', 'iso-8859-1', 'cp1252']

# Constants for text processing
WORD_PATTERN = re.compile(r'\b[a-zA-Z]+\b')  # Matches whole words containing only letters

def get_time_filter():
    """Ask user if they want to filter files by modification time"""
    print("\nTime Filter Configuration")
    print("=======================")
    print("Would you like to scan only files modified after a certain date/time?")
    print("1. No time filter (scan all files)")
    print("2. Last 24 hours")
    print("3. Last 7 days")
    print("4. Last 30 days")
    print("5. Custom date")
    
    while True:
        choice = input("\nYour selection (1-5): ").strip()
        
        current_time = time.time()
        if choice == "1":
            config['use_mtime_filter'] = False
            config['min_mtime'] = None
            print("\nTime filter disabled - scanning all files regardless of modification time.")
            return
        elif choice == "2":
            config['min_mtime'] = current_time - (24 * 3600)  # 24 hours in seconds
            filter_desc = "last 24 hours"
        elif choice == "3":
            config['min_mtime'] = current_time - (7 * 24 * 3600)  # 7 days in seconds
            filter_desc = "last 7 days"
        elif choice == "4":
            config['min_mtime'] = current_time - (30 * 24 * 3600)  # 30 days in seconds
            filter_desc = "last 30 days"
        elif choice == "5":
            while True:
                print("\nEnter date (Format: YYYY-MM-DD)")
                print("Example: 2024-03-15 for March 15, 2024")
                date_str = input("Your date: ").strip()
                try:
                    custom_date = datetime.strptime(date_str, "%Y-%m-%d")
                    config['min_mtime'] = custom_date.timestamp()
                    filter_desc = f"after {date_str}"
                    break
                except ValueError:
                    print("\nInvalid date format. Please use YYYY-MM-DD format.")
                    continue
        else:
            print("Invalid selection. Please enter a number between 1 and 5.")
            continue
        
        config['use_mtime_filter'] = True
        print(f"\nTime filter enabled - only scanning files modified {filter_desc}")
        print(f"Filter time set to: {datetime.fromtimestamp(config['min_mtime']).strftime('%Y-%m-%d %H:%M:%S')}")
        break

def normalize_word(word: str) -> str:
    """Normalize a word by converting to lowercase and removing special characters"""
    try:
        # Convert to lowercase and strip whitespace
        word = word.lower().strip()
        # Remove any non-letter characters
        word = ''.join(c for c in word if c.isalpha())
        return word if word else ''
    except Exception:
        return ''

def assess_seed_quality(phrase: str) -> str:
    """
    Assess the quality of a potential seed phrase.
    Returns 'high' or 'low'.
    """
    try:
        words = phrase.lower().split()
        
        # Basic validation
        if len(words) not in [12, 18, 24]:
            return "low"
            
        if not all(word in BIP39_WORDS for word in words):
            return "low"
            
        if len(set(words)) != len(words):
            return "low"
        
        # Additional quality checks
        quality_score = 0
        
        # 1. Check word distribution
        # Real seed phrases typically have a good mix of word lengths
        word_lengths = [len(word) for word in words]
        length_variety = len(set(word_lengths))
        if length_variety >= 4:  # Good variety of word lengths
            quality_score += 1
            
        # 2. Check for common patterns in real seeds
        # Real seeds rarely have many common English words
        common_words = {'the', 'and', 'that', 'have', 'for', 'not', 'with'}
        common_word_count = sum(1 for word in words if word in common_words)
        if common_word_count <= 1:  # Allow max 1 common word
            quality_score += 1
            
        # 3. Check word uniqueness
        unique_words = len(set(words))
        if unique_words == len(words):  # All words should be unique
            quality_score += 1
            
        # 4. Check for sequential patterns
        has_sequential = False
        for i in range(len(words) - 1):
            if words[i] == words[i + 1]:
                has_sequential = True
                break
        if not has_sequential:
            quality_score += 1
            
        # Require a high quality score to be considered high quality
        return "high" if quality_score >= 3 else "low"
        
    except Exception:
        return "low"

def log_seed(phrase: str, file_path: str, quality: str) -> None:
    """Log a found seed phrase with timestamp and source"""
    try:
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        output_dir = 'output'
        os.makedirs(output_dir, exist_ok=True)
        
        # Use quality level as filename prefix
        output_file = os.path.join(output_dir, f'{quality}_quality_seeds.txt')
        
        # Format the output with more details
        word_count = len(phrase.split())
        output_line = (
            f"{'='*80}\n"
            f"TIMESTAMP: {timestamp}\n"
            f"QUALITY: {quality.upper()}\n"
            f"WORD COUNT: {word_count}\n"
            f"SOURCE FILE: {file_path}\n"
            f"SEED PHRASE: {phrase}\n"
            f"{'='*80}\n\n"
        )
        
        with open(output_file, 'a', encoding='utf-8') as f:
            f.write(output_line)
            
        # For high quality seeds, also save to a separate daily file
        if quality == "high":
            daily_dir = os.path.join(output_dir, 'daily_finds')
            os.makedirs(daily_dir, exist_ok=True)
            daily_file = os.path.join(daily_dir, f"high_quality_seeds_{datetime.now().strftime('%Y%m%d')}.txt")
            
            with open(daily_file, 'a', encoding='utf-8') as f:
                f.write(output_line)
            
            # Print immediate notification for high quality seeds
            print("\n" + "!"*80)
            print(f"HIGH QUALITY {word_count}-WORD SEED PHRASE FOUND!")
            print(f"Location: {file_path}")
            print(f"Phrase: {phrase}")
            print(f"Saved to: {output_file}")
            print("!"*80 + "\n")
            sys.stdout.flush()
            
        stats.seeds_found += 1
        if quality == "high":
            stats.high_quality_seeds += 1
        
    except Exception as e:
        logger.error(f"Error logging seed phrase: {str(e)}")

def find_seed_phrases(text: str) -> List[str]:
    """Find potential seed phrases with enhanced pattern matching"""
    try:
        # Get all potential words
        words = WORD_PATTERN.findall(text.lower())  # Convert text to lowercase first
        
        # Filter out empty words and normalize
        words = [w for w in words if w and w in BIP39_WORDS]
        
        # If we don't have enough BIP39 words, return early
        if len(words) < 12:
            return []
        
        found = []
        # Look for standard BIP39 lengths (12, 18, or 24 words)
        for count in [12, 18, 24]:
            for i in range(len(words) - count + 1):
                chunk = words[i:i + count]
                # Only consider chunks where ALL words are in BIP39
                if all(word in BIP39_WORDS for word in chunk):
                    phrase = ' '.join(chunk)
                    # Additional validation to reduce false positives
                    if is_potential_seed_phrase(chunk):
                        found.append(phrase)
        
        return found
    except Exception as e:
        logger.error(f"Error in find_seed_phrases: {str(e)}")
        return []

def is_potential_seed_phrase(words: List[str]) -> bool:
    """
    Apply additional validation rules to reduce false positives.
    Returns True if the phrase looks like a potential seed phrase.
    """
    # Must be standard BIP39 length
    if len(words) not in [12, 18, 24]:
        return False
        
    # No duplicate words allowed in seed phrases
    if len(set(words)) != len(words):
        return False
        
    # Check for common words that might indicate natural text
    common_text_indicators = {
        'the', 'and', 'that', 'have', 'for', 'not', 'with', 'you', 'this', 'but',
        'his', 'her', 'they', 'will', 'one', 'all', 'would', 'there', 'their'
    }
    
    # If too many common words, probably not a seed phrase
    common_word_count = sum(1 for word in words if word in common_text_indicators)
    if common_word_count > 2:  # Allow max 2 common words
        return False
    
    # Check for sequential words that might indicate natural text
    for i in range(len(words) - 1):
        if words[i] == words[i + 1]:  # No repeated words
            return False
    
    return True

# Target file extensions for scanning inside archives
target_extensions = {
    '.txt', '.log', '.ini', '.conf', '.cfg',
    '.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.php',
    '.html', '.htm', '.css',
    '.json', '.xml', '.csv',
    '.md', '.markdown',
    '.doc', '.docx', '.rtf', '.odt',
    '.xls', '.xlsx', '.ods',
    '.ppt', '.pptx', '.odp',
    '.pdf',
    '.db', '.sqlite', '.sqlite3'
}

# Size limits for different file types (in bytes)
SIZE_LIMITS = {
    'normal': 100 * 1024 * 1024,      # 100MB for normal files
    'compressed': 50 * 1024 * 1024,    # 50MB for compressed files
    'pdf': 10 * 1024 * 1024,          # 10MB for PDF files
    'email': 10 * 1024 * 1024         # 10MB for email files
}

# Global variables
WORDLIST = None  # Global wordlist
BIP39_WORDS = None  # Global BIP39 words

# Statistics tracking
class ScanStats:
    def __init__(self):
        self.start_time = time.time()
        self.files_processed = 0
        self.total_files = 0
        self.bytes_processed = 0
        self.total_bytes = 0
        self.seeds_found = 0
        self.protected_files = 0
        self.unreadable_files = 0
        self.current_file = ""
        self.last_update = time.time()
        self.update_interval = 0.5  # Update every 0.5 seconds
        
        # New statistics
        self.files_skipped = 0
        self.dirs_scanned = 0
        self.file_types_count = {}  # Dictionary to track file types
        self.inaccessible_locations = set()  # Set to track inaccessible locations
        self.inaccessible_log_file = 'inaccessible_locations.txt'  # Single log file
        self.high_prob_locations_file = 'high_probability_locations.txt'  # File for high probability locations
        self.high_quality_seeds = 0  # Counter for high quality seeds
        self.time_filtered_files = 0  # Counter for files skipped due to time filter
        
    def update(self, bytes_read: int = 0, file_completed: bool = False, 
              seed_found: bool = False, protected_file: bool = False,
              file_type: str = None, skipped: bool = False):
        """Update statistics and display if enough time has passed"""
        self.bytes_processed += bytes_read
        if file_completed:
            self.files_processed += 1
            if file_type:
                self.file_types_count[file_type] = self.file_types_count.get(file_type, 0) + 1
        if seed_found:
            self.seeds_found += 1
        if protected_file:
            self.protected_files += 1
        if skipped:
            self.files_skipped += 1
            
        current_time = time.time()
        if current_time - self.last_update >= self.update_interval:
            self.display_progress()
            self.last_update = current_time
    
    def add_inaccessible_location(self, path: str, reason: str):
        """Add an inaccessible location with the reason"""
        self.inaccessible_locations.add((path, reason))
        self._save_inaccessible_location(path, reason)
    
    def _save_inaccessible_location(self, path: str, reason: str):
        """Save inaccessible location to a single log file"""
        try:
            with open(self.inaccessible_log_file, 'a', encoding='utf-8') as f:
                timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                f.write(f"{timestamp} | {path} | {reason}\n")
        except Exception as e:
            logger.error(f"Error saving inaccessible location: {str(e)}")
    
    def display_progress(self):
        """Display real-time progress and statistics"""
        elapsed = time.time() - self.start_time
        if elapsed == 0:
            elapsed = 0.1  # Avoid division by zero
            
        bytes_per_sec = self.bytes_processed / elapsed
        files_per_sec = self.files_processed / elapsed
        
        # Calculate estimated time remaining
        if self.total_bytes > 0:
            bytes_remaining = self.total_bytes - self.bytes_processed
            eta_seconds = bytes_remaining / bytes_per_sec if bytes_per_sec > 0 else 0
        else:
            eta_seconds = 0
            
        # Calculate progress percentage
        progress_pct = (self.files_processed/self.total_files*100) if self.total_files > 0 else 0
            
        # Clear previous lines
        print("\033[K", end="\r")  # Clear current line
        sys.stdout.flush()  # Ensure line clear is processed
        
        # Display current progress
        if self.current_file:
            print(f"Current file: {self.current_file}")
        else:
            print("Initializing scan...")
        sys.stdout.flush()  # Ensure current file is displayed
        
        print(f"Progress: {self.files_processed}/{self.total_files} files ({progress_pct:.1f}%)")
        print(f"Bytes processed: {self.bytes_processed/1024/1024:.1f}MB / "
              f"{self.total_bytes/1024/1024:.1f}MB")
        print(f"Speed: {bytes_per_sec/1024/1024:.1f}MB/s, "
              f"{files_per_sec:.1f} files/s")
        print(f"Seeds found: {self.seeds_found}")
        print(f"Protected files: {self.protected_files}")
        print(f"Files skipped: {self.files_skipped}")
        print(f"ETA: {timedelta(seconds=int(eta_seconds))}")
        sys.stdout.flush()  # Ensure all stats are displayed
        
        # Move cursor up 8 lines to overwrite next update
        print("\033[3A", end="")
        sys.stdout.flush()  # Ensure cursor position is updated
    
    def finalize(self):
        """Display final statistics"""
        elapsed = time.time() - self.start_time
        print("\n" * 8)  # Clear previous progress display
        
        print("\nScan Complete!")
        print("==============")
        print(f"\nTime Statistics:")
        print(f"- Total time: {timedelta(seconds=int(elapsed))}")
        
        print(f"\nFile Statistics:")
        print(f"- Total files found: {self.total_files:,}")
        print(f"- Files processed: {self.files_processed:,}")
        print(f"- Files skipped: {self.files_skipped:,}")
        if config['use_mtime_filter']:
            print(f"  - Skipped due to time filter: {self.time_filtered_files:,}")
        print(f"- Protected files found: {self.protected_files:,}")
        print(f"- Unreadable files found: {self.unreadable_files:,}")
        print(f"- Directories scanned: {self.dirs_scanned:,}")
        print(f"- Total data processed: {self.bytes_processed/1024/1024:.1f}MB")
        print(f"- Average speed: {self.bytes_processed/elapsed/1024/1024:.1f}MB/s")
        
        print(f"\nScan Results:")
        print(f"- Total seeds found: {self.seeds_found:,}")
        print(f"- High quality seeds: {self.high_quality_seeds:,}")
        
        if self.high_quality_seeds > 0:
            print("\nHigh quality seeds have been saved to:")
            print(f"1. output/high_quality_seeds.txt")
            print(f"2. output/daily_finds/high_quality_seeds_{datetime.now().strftime('%Y%m%d')}.txt")
        
        if self.protected_files > 0:
            print(f"\nProtected files have been logged to: protected_files.txt")
        if self.unreadable_files > 0:
            print(f"Unreadable files have been logged to: unreadable_files.txt")
        if self.inaccessible_locations:
            print(f"Inaccessible locations have been logged to: {self.inaccessible_log_file}")

# Global statistics object
stats = ScanStats()

# Create required directories before setting up logging
try:
    os.makedirs('scan_logs', exist_ok=True)
    os.makedirs('output', exist_ok=True)
except Exception as e:
    print(f"\nError: Could not create required directories: {str(e)}")
    sys.exit(1)

# Files to exclude from scanning
output_files = {
    'english.txt',  # BIP39 wordlist
    'wordlist.txt',  # Custom wordlist
    'scan_history.db',  # Database file
    'protected_files.txt',  # List of protected files
    'unreadable_files.txt',  # List of unreadable files
    'inaccessible_locations.txt'  # List of inaccessible locations
}

def setup_logger():
    """Configure logging with both file and console handlers"""
    # Create a unique log file name with timestamp
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    log_file = os.path.join('scan_logs', f'scan_{timestamp}.log')
    
    # Configure the root logger
    logger = logging.getLogger()
    logger.setLevel(logging.INFO)
    
    # Remove any existing handlers
    for handler in logger.handlers[:]:
        logger.removeHandler(handler)
    
    # Create file handler
    file_handler = logging.FileHandler(log_file, encoding='utf-8')
    file_handler.setLevel(logging.DEBUG)
    file_format = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    file_handler.setFormatter(file_format)
    
    # Create console handler
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.WARNING)  # Only warnings and errors to console
    console_format = logging.Formatter('%(levelname)s: %(message)s')
    console_handler.setFormatter(console_format)
    
    # Add handlers to logger
    logger.addHandler(file_handler)
    logger.addHandler(console_handler)
    
    return logger

# Initialize logger globally - will be used throughout the script
logger = setup_logger()

# Required packages for all platforms
required_packages = {
    'core': [
        'python-magic',  # For better MIME type detection
        'chardet',       # For text encoding detection
        'python-docx',   # For DOCX files
        'openpyxl',      # For Excel files
        'python-pptx',   # For PowerPoint files
        'textract',      # For text extraction
        'PyMuPDF',      # For PDF files
        'tqdm',         # For progress bars
        'beautifulsoup4' # For HTML parsing
    ],
    'windows': [
        'pywin32',      # For Windows-specific functionality and DOC file reading
        'python-magic-bin'  # Windows-compatible magic library
    ]
}

# Try to import Windows-specific modules if on Windows
WORD_AVAILABLE = False
if platform.system() == 'Windows':
    try:
        import win32com.client
        WORD_AVAILABLE = True
    except ImportError:
        WORD_AVAILABLE = False

def validate_environment() -> bool:
    """
    Validate the environment before starting.
    Checks for:
    1. Write permissions
    2. Required directories
    
    Returns True if environment is valid, False otherwise.
    """
    logger.info("Starting environment validation...")
    
    # Print dependency requirements
    print("\nRequired Dependencies:")
    print("=====================")
    print("Core packages:")
    for pkg in required_packages['core']:
        print(f"- {pkg}")
    
    if platform.system() == 'Windows':
        print("\nWindows-specific packages:")
        for pkg in required_packages['windows']:
            print(f"- {pkg}")
    
    print("\nTo install all required packages, run:")
    deps = required_packages['core']
    if platform.system() == 'Windows':
        deps.extend(required_packages['windows'])
    print(f"pip install {' '.join(deps)}")
    
    # Check for write permissions in current directory
    try:
        test_file = "write_test.tmp"
        with open(test_file, 'w') as f:
            f.write("test")
        os.remove(test_file)
        logger.debug("Write permission test passed")
    except Exception as e:
        print("\nError: No write permission in current directory")
        print(f"Error details: {str(e)}")
        logger.error(f"Write permission test failed: {str(e)}")
        return False
    
    # Create required directories if they don't exist
    required_dirs = ["scan_logs", "output"]
    for directory in required_dirs:
        try:
            os.makedirs(directory, exist_ok=True)
            logger.debug(f"Successfully created/verified directory: {directory}")
        except Exception as e:
            print(f"\nError: Could not create required directory '{directory}'")
            print(f"Error details: {str(e)}")
            logger.error(f"Failed to create directory {directory}: {str(e)}")
            return False
    
    logger.info("Environment validation completed successfully")
    return True

def check_dependencies() -> Dict[str, str]:
    """
    Check if all required packages are installed.
    Returns a dictionary of missing packages and their import errors.
    """
    missing = {}
    
    # Check core packages
    for package in required_packages['core']:
        try:
            __import__(package.replace('-', '_'))
        except ImportError as e:
            missing[package] = str(e)
    
    # Check platform-specific packages
    if platform.system() == 'Windows':
        for package in required_packages['windows']:
            try:
                __import__(package.replace('-', '_'))
            except ImportError as e:
                missing[package] = str(e)
    
    return missing

# Dictionary to store import errors
missing_dependencies = {}

# Try importing external dependencies
try:
    import fitz  # PyMuPDF
except ImportError as e:
    missing_dependencies['PyMuPDF'] = str(e)

try:
    import docx
except ImportError as e:
    missing_dependencies['python-docx'] = str(e)

try:
    from bs4 import BeautifulSoup
except ImportError as e:
    missing_dependencies['beautifulsoup4'] = str(e)

# Windows-specific imports are now conditional
windows_imports_available = False
if platform.system() == 'Windows':
    try:
        import win32api
        import win32file
        import win32con
        windows_imports_available = True
    except ImportError as e:
        missing_dependencies['pywin32'] = str(e)

# Memory management - limit seen phrases cache
MAX_SEEN_PHRASES = 500000  # Limit cache to 500k phrases
seen_phrases: Set[str] = set()

def clear_phrase_cache():
    """Clear the phrase cache if it gets too large"""
    global seen_phrases
    if len(seen_phrases) > MAX_SEEN_PHRASES:
        seen_phrases.clear()

def get_os_selection():
    """Get user's operating system selection"""
    system = platform.system()
    while True:
        print("\nOperating System Selection")
        print("========================")
        print("Please select your operating system:")
        print("1. Windows")
        print("2. macOS")
        print("3. Linux")
        print(f"\nDetected OS: {system}")
        print("\nType the number (1-3) that corresponds to your operating system.")
        choice = input("Your selection: ").strip()
        
        if choice == "1":
            if system != 'Windows':
                print("\nWarning: You selected Windows but we detected a different OS.")
                confirm = input("Are you sure you want to proceed with Windows? (y/n): ").lower()
                if confirm != 'y':
                    continue
            if not windows_imports_available:
                print("\nError: Windows-specific dependencies (pywin32) are required but not available.")
                print("Please install pywin32 using: pip install pywin32")
                sys.exit(1)
            config['drive_letters'] = True  # Enable drive letters for Windows
            return "windows"
        elif choice == "2":
            if system != 'Darwin':
                print("\nWarning: You selected macOS but we detected a different OS.")
                confirm = input("Are you sure you want to proceed with macOS? (y/n): ").lower()
                if confirm != 'y':
                    continue
            config['drive_letters'] = False  # Disable drive letters for macOS
            return "macos"
        elif choice == "3":
            if system != 'Linux':
                print("\nWarning: You selected Linux but we detected a different OS.")
                confirm = input("Are you sure you want to proceed with Linux? (y/n): ").lower()
                if confirm != 'y':
                    continue
            config['drive_letters'] = False  # Disable drive letters for Linux
            return "linux"
        else:
            print("\nInvalid selection. Please enter 1, 2, or 3.")

def safe_path_exists(path: str) -> bool:
    """Safely check if a path exists handling various edge cases"""
    try:
        return Path(path).exists() and os.access(path, os.R_OK)
    except (OSError, ValueError, TypeError):
        return False

def safe_path_join(*args: str) -> str:
    """Safely join path components handling various edge cases"""
    try:
        return str(Path(*args).resolve())
    except (TypeError, ValueError, OSError):
        return ""

def safe_get_size(path: str) -> int:
    """Safely get file size handling various edge cases"""
    try:
        return Path(path).stat().st_size if Path(path).is_file() else 0
    except (OSError, ValueError, TypeError):
        return 0

def get_network_drives():
    """Get all mapped network drives based on the operating system"""
    network_drives = []
    
    if OPERATING_SYSTEM == "windows":
        try:
            if windows_imports_available:
                drives = win32api.GetLogicalDriveStrings().split('\000')[:-1]
                for drive in drives:
                    try:
                        if win32file.GetDriveType(drive) == win32con.DRIVE_REMOTE:
                            network_drives.append(drive)
                    except Exception:
                        continue
            else:
                # Fallback method for Windows without pywin32
                for letter in string.ascii_uppercase:
                    path = f"{letter}:\\"
                    if safe_path_exists(path):
                        try:
                            if os.path.realpath(path).startswith('\\\\'):
                                network_drives.append(path)
                        except Exception:
                            continue
        except Exception as e:
            logger.error(f"Error detecting network drives: {str(e)}")
    
    elif OPERATING_SYSTEM == "macos":
        volumes_dir = "/Volumes"
        if safe_path_exists(volumes_dir):
            try:
                for volume in os.listdir(volumes_dir):
                    path = safe_path_join(volumes_dir, volume)
                    if os.path.ismount(path) and volume != "Macintosh HD":
                        network_drives.append(path)
            except Exception as e:
                logger.error(f"Error detecting network drives: {str(e)}")
    
    elif OPERATING_SYSTEM == "linux":
        mount_points = ["/mnt", "/media/" + os.getenv("USER", "")]
        for mount_point in mount_points:
            if safe_path_exists(mount_point):
                try:
                    for item in os.listdir(mount_point):
                        path = safe_path_join(mount_point, item)
                        if os.path.ismount(path):
                            network_drives.append(path)
                except Exception as e:
                    logger.error(f"Error detecting network drives: {str(e)}")
    
    return network_drives

def load_wordlist() -> Set[str]:
    """
    Load words from either wordlist.txt if it exists in the script directory,
    or fall back to english.txt for the default BIP39 wordlist.
    """
    def display_and_confirm_wordlist(words: Set[str], list_type: str) -> bool:
        """Helper function to display and confirm wordlist"""
        word_count = len(words)
        sorted_words = sorted(words)
        
        print(f"\n{list_type} Wordlist Contents")
        print("=" * (len(list_type) + 17))
        print(f"\nTotal words: {word_count:,}")
        print("\nFirst 20 words:")
        for i, word in enumerate(sorted_words[:20], 1):
            print(f"{i:2d}. {word}")
        
        if word_count > 20:
            print("\n...")  # Indicate there are more words
            
            print("\nLast 20 words:")
            for i, word in enumerate(sorted_words[-20:], word_count-19):
                print(f"{i:2d}. {word}")
        
        while True:
            confirm = input(f"\nIs this the correct {list_type.lower()} wordlist you want to use? (yes/no): ").lower()
            if confirm == 'no':
                return False
            elif confirm == 'yes':
                return True
            else:
                print("Please enter 'yes' or 'no'")
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    custom_wordlist = os.path.join(script_dir, 'wordlist.txt')
    english_wordlist = os.path.join(script_dir, 'english.txt')
    
    # First try custom wordlist if it exists
    if os.path.exists(custom_wordlist):
        try:
            with open(custom_wordlist, 'r', encoding='utf-8') as f:
                words = {word.strip().lower() for word in f if word.strip()}
            if words:  # Only use if file contains words
                if display_and_confirm_wordlist(words, "Custom"):
                    return words
        except Exception as e:
            logger.error(f"Error loading wordlist.txt: {str(e)}")
    
    # Load and confirm english.txt (BIP39 wordlist)
    if not os.path.exists(english_wordlist):
        print("\nError: Could not find english.txt (BIP39 wordlist)")
        print("Please ensure english.txt is in the same directory as this script.")
        print("You can download it from: https://github.com/bitcoin/bips/blob/master/bip-0039/english.txt")
        sys.exit(1)
        
    try:
        with open(english_wordlist, 'r', encoding='utf-8') as f:
            words = {word.strip().lower() for word in f if word.strip()}
        if not words:
            print("\nError: english.txt is empty")
            sys.exit(1)
        print("\nUsing BIP39 wordlist from english.txt")
        if display_and_confirm_wordlist(words, "BIP39"):
            return words
        else:
            print("\nBIP39 wordlist rejected. Please check english.txt or provide a custom wordlist.")
            sys.exit(1)
    except Exception as e:
        print(f"\nError reading english.txt: {str(e)}")
        print("Please ensure english.txt is a valid UTF-8 encoded file.")
        sys.exit(1)

def scan_doc(file_path: str) -> None:
    """
    Scan legacy DOC files using multiple methods:
    1. Try win32com (if on Windows and Word is available)
    2. Try direct binary reading with multiple encodings
    """
    try:
        # First try win32com on Windows if Word is available
        if platform.system() == 'Windows' and WORD_AVAILABLE:
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                try:
                    doc = word.Documents.Open(file_path)
                    text = doc.Content.Text
                    doc.Close()
                    word.Quit()
                    if text.strip():
                        scan_text(text, file_path)
                        logger.info(f"Successfully extracted text from {os.path.basename(file_path)} using Word")
                        return
                except Exception as e:
                    logger.warning(f"Could not read {os.path.basename(file_path)} with Word: {str(e)}")
                finally:
                    try:
                        word.Quit()
                    except:
                        pass
            except Exception as e:
                logger.warning(f"Microsoft Word automation failed for {os.path.basename(file_path)}")

        # Try direct binary reading as fallback
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
                
                # Try different encodings
                text = None
                for encoding in ['utf-16le', 'utf-8', 'ascii', 'cp1252']:
                    try:
                        decoded = content.decode(encoding, errors='ignore')
                        if len(decoded.strip()) > 0:
                            text = decoded
                            break
                    except Exception:
                        continue
                
                if text:
                    # Clean up the text
                    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)  # Remove control chars
                    text = re.sub(r'\s+', ' ', text)  # Normalize whitespace
                    text = text.strip()
                    
                    if len(text) > 0:
                        scan_text(text, file_path)
                        logger.info(f"Successfully extracted text from {os.path.basename(file_path)} using binary reading")
                    else:
                        logger.warning(f"No readable text found in {os.path.basename(file_path)}")
                else:
                    logger.warning(f"Could not decode {os.path.basename(file_path)} with any supported encoding")
                    
        except Exception as e:
            logger.error(f"Error reading {os.path.basename(file_path)}: {str(e)}")
            
    except Exception as e:
        logger.error(f"Error processing {os.path.basename(file_path)}: {str(e)}")
        if platform.system() == 'Windows':
            logger.info("\nFor better DOC file support on Windows:")
            logger.info("1. Install Microsoft Word")
            logger.info("2. Or convert DOC files to DOCX using online converters")
            if not WORD_AVAILABLE:
                logger.info("3. Install pywin32: pip install pywin32")
        else:
            logger.info("\nFor better DOC file support on Linux/macOS:")
            logger.info("- Linux: sudo apt-get install catdoc")
            logger.info("- macOS: brew install catdoc")

def scan_text(text: str, file_path: str) -> None:
    """Scan text for potential seed phrases"""
    try:
        matches = find_seed_phrases(text)
        for match in matches:
            # Only log if we haven't seen this exact phrase before
            if match not in seen_phrases:
                quality = assess_seed_quality(match)
                # Log all potential matches but mark quality
                log_seed(match, file_path, quality)
                seen_phrases.add(match)
                
                # For all potential matches (even low quality), print to console
                print("\n" + "-"*80)
                print(f"Potential seed phrase found in: {os.path.basename(file_path)}")
                print(f"Quality: {quality.upper()}")
                print(f"Words: {match}")
                print("-"*80 + "\n")
                sys.stdout.flush()
                
                if quality == "high":
                    logger.info(f"Found high quality seed phrase in {file_path}")
                else:
                    logger.debug(f"Found potential low quality seed phrase in {file_path}")
    except Exception as e:
        logger.error(f"Error scanning text from {file_path}: {str(e)}")

def scan_text_file(file_path: str) -> List[str]:
    """
    Scan a text file for potential seed phrases, handling different encodings.
    Returns a list of found phrases.
    """
    found_phrases = []
    
    # Try different encodings
    for encoding in ENCODINGS:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
                # Process the content
                words = WORD_PATTERN.findall(content.lower())
                
                # Check for consecutive words that match the wordlist
                for i in range(len(words) - 11):  # Look for 12-word phrases
                    phrase = words[i:i+12]
                    if all(word in BIP39_WORDS for word in phrase):
                        found_phrases.append(' '.join(phrase))
                
                # If we successfully read the file, break the loop
                break
                
        except UnicodeDecodeError:
            continue
        except Exception as e:
            logger.debug(f"Error reading {file_path} with {encoding} encoding: {str(e)}")
            continue
    
    # If all encodings failed, try binary mode with UTF-8 and ignore errors
    if not found_phrases:
        try:
            with open(file_path, 'rb') as f:
                content = f.read().decode('utf-8', errors='ignore')
                words = WORD_PATTERN.findall(content.lower())
                
                for i in range(len(words) - 11):
                    phrase = words[i:i+12]
                    if all(word in BIP39_WORDS for word in phrase):
                        found_phrases.append(' '.join(phrase))
                        
        except Exception as e:
            logger.error(f"Failed to read {file_path} in binary mode: {str(e)}")
    
    return found_phrases

def safe_scan_file(file_path: str) -> None:
    """Safely scan a file with comprehensive error handling and optimized type detection"""
    try:
        if not safe_path_exists(file_path):
            logger.debug(f"File not accessible: {file_path}")
            return

        filename = os.path.basename(file_path)
        if filename in output_files:
            logger.debug(f"Skipping output file to prevent recursive scanning: {filename}")
            return

        # Log high probability locations
        if is_high_probability_directory(file_path):
            log_high_probability_location(file_path)

        try:
            file_type, is_protected = get_file_type(file_path)
            
            if is_protected:
                log_protected_file(file_path)
                return
                
            if file_type == 'pdf':
                scan_pdf(file_path)
            elif file_type == 'document':
                scan_document(file_path)
            elif file_type == 'spreadsheet':
                scan_spreadsheet(file_path)
            elif file_type == 'presentation':
                scan_presentation(file_path)
            elif file_type == 'database':
                scan_database(file_path)
            elif file_type == 'compressed':
                scan_zip(file_path)
            elif file_type == 'text':
                scan_text_file(file_path)
            else:
                logger.debug(f"Skipping unsupported file type: {file_path}")
        except MemoryError:
            logger.error(f"Out of memory while processing {file_path}")
            return
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            return
            
    except Exception as e:
        logger.error(f"Error accessing file {file_path}: {str(e)}")

def scan_pdf(file_path: str) -> None:
    """Scan a PDF file with improved error handling"""
    if not PDF_SUPPORT:
        logger.error(f"PDF support not available. Please install PyMuPDF: pip install PyMuPDF")
        return

    doc = None
    try:
        doc = fitz.open(file_path)
        page_count = len(doc)
        
        if config['exclude_pdf_books'] and page_count > config['pdf_book_pages']:
            logger.debug(f"Skipping PDF book ({page_count} pages): {os.path.basename(file_path)}")
            return
            
        text = ""
        for page_num in range(page_count):
            try:
                page = doc[page_num]
                text += page.get_text() + "\n"
            except Exception as e:
                logger.warning(f"Error reading page {page_num + 1} in {os.path.basename(file_path)}: {str(e)}")
                continue

        if text.strip():
            scan_text(text, file_path)
            logger.info(f"Successfully scanned PDF: {os.path.basename(file_path)}")
        else:
            logger.warning(f"No readable text found in PDF: {os.path.basename(file_path)}")

    except Exception as e:
        logger.error(f"Error processing PDF {os.path.basename(file_path)}: {str(e)}")
        logger.info("\nFor PDF support, please ensure PyMuPDF is installed:")
        logger.info("pip install PyMuPDF")
    finally:
        if doc:
            try:
                doc.close()
            except Exception:
                pass

def scan_docx(file_path: str) -> None:
    """Scan a DOCX file"""
    if not DOCX_SUPPORT:
        logger.error(f"DOCX support not available. Please install python-docx: pip install python-docx")
        return

    try:
        doc = Document(file_path)
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        
        # Also check tables if present
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += "\n" + cell.text

        if text.strip():
            scan_text(text, file_path)
            logger.info(f"Successfully scanned DOCX: {os.path.basename(file_path)}")
        else:
            logger.warning(f"No readable text found in DOCX: {os.path.basename(file_path)}")
            
    except Exception as e:
        logger.error(f"Error processing DOCX {os.path.basename(file_path)}: {str(e)}")
        logger.info("\nFor DOCX support, please ensure python-docx is installed:")
        logger.info("pip install python-docx")

def scan_odt(file_path: str) -> None:
    """Scan OpenDocument Text files"""
    try:
        # ODT files are ZIP files containing content.xml
        with zipfile.ZipFile(file_path) as odt:
            with odt.open('content.xml') as content:
                text = content.read().decode('utf-8')
                scan_text(text, file_path)
    except Exception as e:
        logger.error(f"Error scanning ODT file {file_path}: {str(e)}")

def scan_rtf(file_path: str) -> None:
    """Scan RTF files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            # Strip RTF formatting
            text = re.sub(r'[\\\{\}]|\\\w+\s?', '', content)
            scan_text(text, file_path)
    except Exception as e:
        logger.error(f"Error scanning RTF file {file_path}: {str(e)}")

def scan_document(file_path: str) -> None:
    """Scan various document formats (doc, docx, odt, rtf)"""
    try:
        file_type, is_protected = get_file_type(file_path)
        if is_protected:
            log_protected_file(file_path)
            return
            
        if file_path.endswith('.docx'):
            scan_docx(file_path)
        elif file_path.endswith('.odt'):
            scan_odt(file_path)
        elif file_path.endswith('.rtf'):
            scan_rtf(file_path)
        elif file_path.endswith('.doc'):
            scan_doc(file_path)
    except Exception as e:
        logger.error(f"Error scanning document {file_path}: {str(e)}")

def scan_spreadsheet(file_path: str) -> None:
    """Scan spreadsheet formats (xls, xlsx, ods)"""
    try:
        file_type, is_protected = get_file_type(file_path)
        if is_protected:
            log_protected_file(file_path)
            return
            
        # Extract text content from cells
        if file_path.endswith(('.xlsx', '.xls')):
            scan_excel(file_path)
        elif file_path.endswith('.ods'):
            scan_ods(file_path)
    except Exception as e:
        logger.error(f"Error scanning spreadsheet {file_path}: {str(e)}")

def scan_excel(file_path: str) -> None:
    """Scan Excel files (XLS/XLSX)"""
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb:
            text = []
            for row in sheet.iter_rows():
                text.extend(str(cell.value) for cell in row if cell.value)
            scan_text('\n'.join(text), file_path)
    except Exception as e:
        logger.error(f"Error scanning Excel file {file_path}: {str(e)}")

def scan_ods(file_path: str) -> None:
    """Scan OpenDocument Spreadsheet files"""
    try:
        text = textract.process(file_path).decode('utf-8')
        scan_text(text, file_path)
    except Exception as e:
        logger.error(f"Error scanning ODS file {file_path}: {str(e)}")

def scan_presentation(file_path: str) -> None:
    """Scan presentation formats (ppt, pptx, odp)"""
    try:
        file_type, is_protected = get_file_type(file_path)
        if is_protected:
            log_protected_file(file_path)
            return
            
        if file_path.endswith(('.pptx', '.ppt')):
            scan_powerpoint(file_path)
        elif file_path.endswith('.odp'):
            scan_odp(file_path)
    except Exception as e:
        logger.error(f"Error scanning presentation {file_path}: {str(e)}")

def scan_powerpoint(file_path: str) -> None:
    """Scan PowerPoint files (PPT/PPTX)"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        scan_text("\n".join(text), file_path)
    except Exception as e:
        logger.error(f"Error scanning PowerPoint file {file_path}: {str(e)}")

def scan_odp(file_path: str) -> None:
    """Scan OpenDocument Presentation files"""
    try:
        text = textract.process(file_path).decode('utf-8')
        scan_text(text, file_path)
    except Exception as e:
        logger.error(f"Error scanning ODP file {file_path}: {str(e)}")

# OS-specific configurations
os_config = {
    "windows": {
        "excluded_directories": [
            'C:\\Windows',
            'C:\\Program Files',
            'C:\\Program Files (x86)',
            'C:\\$Recycle.Bin',
            'C:\\System Volume Information'
        ],
        "home_dir": os.path.expanduser('~'),
        "drive_letters": True
    },
    "macos": {
        "excluded_directories": [
            '/System',
            '/Library',
            '/private',
            '.Trash'
        ],
        "home_dir": os.path.expanduser('~'),
        "drive_letters": False
    },
    "linux": {
        "excluded_directories": [
            '/bin',
            '/boot',
            '/dev',
            '/etc',
            '/lib',
            '/lib64',
            '/proc',
            '/sys',
            '/usr/bin',
            '/usr/lib',
            '/var'
        ],
        "home_dir": os.path.expanduser('~'),
        "drive_letters": False
    }
}

def confirm_excluded_directories(os_type):
    """Ask user to confirm which directories to exclude"""
    default_excluded = os_config[os_type]["excluded_directories"]
    confirmed_excluded = []
    
    print("\nExcluded Directories Configuration")
    print("================================")
    print("By default, the following directories are excluded from scanning:")
    for i, directory in enumerate(default_excluded, 1):
        print(f"{i}. {directory}")
    
    while True:
        print("\nWould you like to:")
        print("1. Keep all these directories excluded (recommended)")
        print("2. Include all directories in the scan")
        print("3. Manually select which directories to exclude")
        choice = input("\nYour selection (1-3): ").strip()
        
        if choice == "1":
            return default_excluded
        elif choice == "2":
            print("\nWarning: Including system directories may significantly increase scan time")
            confirm = input("Are you sure you want to include all directories? (y/n): ").lower()
            if confirm == 'y':
                return []
            else:
                continue
        elif choice == "3":
            print("\nFor each directory, type 'y' to exclude it or 'n' to include it in the scan:")
            for directory in default_excluded:
                while True:
                    choice = input(f"Exclude {directory}? (y/n): ").lower()
                    if choice in ['y', 'n']:
                        if choice == 'y':
                            confirmed_excluded.append(directory)
                        break
                    print("Please enter 'y' for yes or 'n' for no.")
            return confirmed_excluded
        else:
            print("Invalid selection. Please enter 1, 2, or 3.")

# Initial configuration
config = {
    'min_file_size': 0,  # scan all files
    'max_file_size': 500 * 1024 * 1024,  # 500MB
    'excluded_directories': [],  # Will be updated based on OS
    'high_probability_directories': [
        'Desktop',
        'Downloads',
        'Documents',
        '.exodus',
        'wallet',
        'blockchain',
        'metamask',
        'crypto'
    ],
    'log_directory': 'scan_logs',
    'debug_mode': False,
    'pdf_book_pages': 50,  # threshold for considering a PDF as a book
    'exclude_pdf_books': True,  # whether to skip PDFs that look like books
    'min_mtime': None,  # minimum modification time (will be set based on user input)
    'use_mtime_filter': False,  # whether to use modification time filter
    'drive_letters': False  # will be updated based on OS selection
}

def should_process_file(db: 'DatabaseManager', file_path: str) -> bool:
    """Determine if a file should be processed"""
    try:
        # Never skip our output files when writing to them
        if os.path.basename(file_path) in output_files:
            return True
            
        try:
            mtime = os.path.getmtime(file_path)
        except OSError as e:
            logger.error(f"Could not get modification time for {file_path}: {str(e)}")
            stats.unreadable_files += 1
            return False
        
        # Check modification time filter if enabled
        if config['use_mtime_filter'] and config['min_mtime'] is not None:
            try:
                if mtime < config['min_mtime']:
                    # Add debug logging
                    logger.debug(f"File: {file_path}")
                    logger.debug(f"File modification time: {datetime.fromtimestamp(mtime).strftime('%Y-%m-%d %H:%M:%S')}")
                    logger.debug(f"Filter minimum time: {datetime.fromtimestamp(config['min_mtime']).strftime('%Y-%m-%d %H:%M:%S')}")
                    stats.time_filtered_files += 1
                    stats.files_skipped += 1
                    return False
            except Exception as e:
                logger.error(f"Error comparing modification times for {file_path}: {str(e)}")
                logger.error(f"mtime: {mtime}, min_mtime: {config['min_mtime']}")
                return False
            
        # File size checks
        try:
            file_size = os.path.getsize(file_path)
        except OSError as e:
            logger.error(f"Could not get size for {file_path}: {str(e)}")
            return False
        
        # Determine appropriate size limit based on file type
        size_limit = SIZE_LIMITS['normal']  # default
        if is_compressed_file(file_path):
            size_limit = SIZE_LIMITS['compressed']
        elif file_path.lower().endswith('.pdf'):
            size_limit = SIZE_LIMITS['pdf']
        elif any(file_path.lower().endswith(ext) for ext in ['.pst', '.ost', '.mbox', '.mbx']):
            size_limit = SIZE_LIMITS['email']
        
        # Skip if file size is outside our bounds
        if file_size < config['min_file_size'] or file_size > size_limit:
            stats.files_skipped += 1
            return False
            
        # Create a simple file identifier using size and mtime
        file_id = f"{file_size}_{mtime}"
        
        # Check if already processed
        if db.was_processed(file_path, file_id):
            logger.debug(f"Skipping previously processed file: {file_path}")
            stats.files_skipped += 1
            return False
        
        # Add to processed files
        db.add_processed_file(file_path, file_id, mtime)
        return True
        
    except Exception as e:
        logger.error(f"Error checking file {file_path}: {str(e)}")
        return False

def ask_resume_scan() -> bool:
    """Ask user if they want to resume previous scan"""
    print("\nChecking for incomplete scans...")
    db = DatabaseManager()
    progress = db.get_last_scan_progress()
    if progress:
        print(f"\nFound incomplete scan:")
        print(f"- Total files: {progress['total_files']:,}")
        print(f"- Files processed: {progress['processed_files']:,}")
        print(f"- Progress: {(progress['processed_files'] / progress['total_files'] * 100):.1f}%")
        
        while True:
            choice = input("\nWould you like to resume this scan? (y/n): ").lower()
            if choice in ['y', 'n']:
                return choice == 'y'
            print("Please enter 'y' for yes or 'n' for no.")
    return False

def is_high_quality_seed(words: List[str]) -> bool:
    """Check if a list of words forms a high quality seed phrase"""
    # All words must be in BIP39 list
    return all(word in BIP39_WORDS for word in words)

def try_read_file(file_path: str) -> Optional[str]:
    """Try reading a file with multiple encoding support and better error handling"""
    if not safe_path_exists(file_path):
        return None
        
    for encoding in ENCODINGS:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
                # Quick validation of content
                content.encode('utf-8')  # Test if content is valid
                return content
        except (UnicodeError, OSError, IOError):
            continue
        except Exception as e:
            logger.debug(f"Unexpected error reading {file_path} with {encoding}: {str(e)}")
            continue
    return None

def get_file_type(file_path: str) -> Tuple[str, bool]:
    """
    Determines file type and whether it's password protected.
    Returns tuple of (file_type, is_protected)
    """
    try:
        # First try extension-based detection for common file types
        ext = os.path.splitext(file_path)[1].lower()
        is_protected = False
        
        # Quick extension check for common types
        if ext in {'.js', '.ts', '.jsx', '.tsx', '.py', '.java', '.cpp', '.c', '.h', '.cs', '.php'}:
            return ('text', False)
        
        # For other files, try python-magic
        try:
            mime_type = magic.from_file(file_path, mime=True)
        except Exception as magic_error:
            logger.debug(f"Magic library error for {file_path}: {str(magic_error)}")
            # Fall back to extension-based detection
            if ext in {'.pdf'}:
                return ('pdf', is_protected)
            elif ext in {'.doc', '.docx', '.odt', '.rtf'}:
                return ('document', is_protected)
            elif ext in {'.xls', '.xlsx', '.ods'}:
                return ('spreadsheet', is_protected)
            elif ext in {'.ppt', '.pptx', '.odp'}:
                return ('presentation', is_protected)
            elif ext in {'.db', '.sqlite', '.sqlite3'}:
                return ('database', is_protected)
            elif ext in {'.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'}:
                return ('compressed', is_protected)
            elif ext in {
                '.txt', '.log', '.ini', '.conf', '.cfg',
                '.html', '.htm', '.css',
                '.json', '.xml', '.csv',
                '.md', '.markdown'
            }:
                return ('text', is_protected)
            return ('unknown', is_protected)
        
        # Check for password protection
        try:
            with open(file_path, 'rb') as f:
                header = f.read(4096)  # Read first 4KB for header analysis
                
                # Common password protection markers
                protection_markers = [
                    b'Encrypted',
                    b'password protected',
                    b'EncryptedPackage',
                    b'Microsoft.Container.EncryptionTransform',
                    b'PDFEncrypted',
                    b'Standard Security'  # PDF encryption marker
                ]
                
                is_protected = any(marker in header for marker in protection_markers)
                
                if is_protected:
                    log_protected_file(file_path, "Encryption detected in file header")
                    
        except Exception as e:
            log_unreadable_file(file_path, f"Cannot read file header: {str(e)}")
            return ('unknown', True)
            
        # Map mime types to our internal types
        if mime_type.startswith('application/pdf'):
            return ('pdf', is_protected)
            
        elif mime_type in {
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # docx
            'application/msword',  # doc
            'application/vnd.oasis.opendocument.text',  # odt
            'application/rtf'  # rtf
        }:
            return ('document', is_protected)
            
        elif mime_type in {
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # xlsx
            'application/vnd.ms-excel',  # xls
            'application/vnd.oasis.opendocument.spreadsheet'  # ods
        }:
            return ('spreadsheet', is_protected)
            
        elif mime_type in {
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # pptx
            'application/vnd.ms-powerpoint',  # ppt
            'application/vnd.oasis.opendocument.presentation'  # odp
        }:
            return ('presentation', is_protected)
            
        elif mime_type == 'application/x-sqlite3' or mime_type.endswith('sqlite'):
            return ('database', is_protected)
            
        elif any(mime_type.startswith(prefix) for prefix in ['application/zip', 'application/x-rar', 'application/x-7z', 'application/x-tar', 'application/gzip']):
            return ('compressed', is_protected)
            
        # Text file detection with encoding check
        elif mime_type.startswith('text/') or mime_type in {
            'application/json',
            'application/xml',
            'application/javascript',
            'application/x-python',
            'application/x-php'
        }:
            return ('text', is_protected)
            
        return ('unknown', is_protected)
        
    except Exception as e:
        logger.debug(f"Error detecting file type for {file_path}: {str(e)}")
        # Fallback to extension-based detection
        ext = os.path.splitext(file_path)[1].lower()
        is_protected = False
        
        if ext in {'.pdf'}:
            return ('pdf', is_protected)
        elif ext in {'.doc', '.docx', '.odt', '.rtf'}:
            return ('document', is_protected)
        elif ext in {'.xls', '.xlsx', '.ods'}:
            return ('spreadsheet', is_protected)
        elif ext in {'.ppt', '.pptx', '.odp'}:
            return ('presentation', is_protected)
        elif ext in {'.db', '.sqlite', '.sqlite3'}:
            return ('database', is_protected)
        elif ext in {'.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'}:
            return ('compressed', is_protected)
        elif ext in {
            '.txt', '.log', '.ini', '.conf', '.cfg',
            '.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.php',
            '.html', '.htm', '.css',
            '.json', '.xml', '.csv',
            '.md', '.markdown'
        }:
            return ('text', is_protected)
            
        return ('unknown', is_protected)

def detect_text_encoding(file_path: str) -> Optional[str]:
    """
    Detects the encoding of a text file.
    Returns the encoding name or None if detection fails.
    """
    try:
        with open(file_path, 'rb') as f:
            raw_data = f.read(4096)  # Read first 4KB for encoding detection
            result = chardet.detect(raw_data)
            if result['confidence'] > 0.7:  # Only accept high confidence results
                return result['encoding']
    except Exception as e:
        logger.error(f"Error detecting encoding for {file_path}: {str(e)}")
    return None

def log_protected_file(file_path: str, reason: str = "Password protected") -> None:
    """
    Logs the location of a protected/encrypted file with timestamp and reason.
    """
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    folder_path = os.path.dirname(file_path)
    filename = os.path.basename(file_path)
    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
    
    protected_files_log = 'protected_files.txt'
    try:
        with open(protected_files_log, 'a', encoding='utf-8') as f:
            f.write(f"{timestamp} | {folder_path} | {filename} | {file_size/1024:.1f}KB | {reason}\n")
        stats.protected_files += 1
    except Exception as e:
        logger.error(f"Error logging protected file {file_path}: {str(e)}")

def log_unreadable_file(file_path: str, error_msg: str) -> None:
    """
    Logs the location of an unreadable file with timestamp and error details.
    """
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    folder_path = os.path.dirname(file_path)
    filename = os.path.basename(file_path)
    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
    
    unreadable_files_log = 'unreadable_files.txt'
    try:
        with open(unreadable_files_log, 'a', encoding='utf-8') as f:
            f.write(f"{timestamp} | {folder_path} | {filename} | {file_size/1024:.1f}KB | {error_msg}\n")
        stats.unreadable_files += 1
    except Exception as e:
        logger.error(f"Error logging unreadable file {file_path}: {str(e)}")

def is_excluded_path(path: str) -> bool:
    """Check if a path should be excluded from scanning"""
    # Convert path to lowercase for case-insensitive comparison
    path_lower = path.lower()
    
    # Check if path is in excluded directories
    for excluded in config['excluded_directories']:
        if path_lower.startswith(excluded.lower()):
            return True
            
    # Check if filename is in output files
    filename = os.path.basename(path)
    if filename in output_files:
        return True
        
    return False

def get_available_drives():
    """Get list of available drives on Windows"""
    if OPERATING_SYSTEM != "windows":
        return []
    
    available_drives = []
    try:
        import string
        from ctypes import windll
        
        # Get bitmask of available drives
        bitmask = windll.kernel32.GetLogicalDrives()
        
        # Check each possible drive letter
        for letter in string.ascii_uppercase:
            if bitmask & (1 << (ord(letter) - ord('A'))):
                drive = f"{letter}:\\"
                if os.path.exists(drive):
                    available_drives.append(drive)
        
        if not available_drives:
            print("\nError: No drives were detected on your system!")
            print("This could indicate a problem with drive access permissions.")
            print("The script will fall back to scanning the current drive only.")
            logger.error("Drive detection failed: No drives found")
            return [os.path.abspath(os.sep)]
            
    except ImportError as e:
        print("\nError: Could not load required Windows components for drive detection!")
        print("This might be because:")
        print("1. The ctypes module is not available")
        print("2. The Windows API is not accessible")
        print("\nFalling back to scanning the current drive only.")
        logger.error(f"Drive detection failed: Import error - {str(e)}")
        return [os.path.abspath(os.sep)]
    except Exception as e:
        print("\nError: Failed to detect system drives!")
        print(f"Error details: {str(e)}")
        print("This could be due to:")
        print("1. Insufficient permissions")
        print("2. System API restrictions")
        print("3. Anti-virus software interference")
        print("\nFalling back to scanning the current drive only.")
        logger.error(f"Drive detection failed: {str(e)}")
        return [os.path.abspath(os.sep)]
    
    return available_drives

def main():
    """Main function with statistics tracking"""
    try:
        # Initialize statistics
        stats.start_time = time.time()
        
        # Get starting paths based on OS
        start_paths = []
        if OPERATING_SYSTEM == "windows":
            if config['drive_letters']:
                print("\nDetecting available drives...")
                # Get all available drives on Windows
                start_paths = get_available_drives()
                if len(start_paths) == 1 and start_paths[0] == os.path.abspath(os.sep):
                    # This means drive detection failed and we're using fallback
                    print("\nWarning: Limited to scanning current drive only.")
                    print("Some drives may be inaccessible.")
                else:
                    print(f"\nSuccessfully detected {len(start_paths)} drive(s)!")
            else:
                start_paths = [config['home_dir']]
        else:
            # For Linux/macOS, start from home directory
            start_paths = [config['home_dir']]
        
        print("\nInitial File System Scan")
        print("=====================")
        print("Starting scan from:")
        for path in start_paths:
            print(f"- {path}")
        print("\nIndexing files and calculating total size...")
        print("(This may take a while for large directories)")
        print("\nPress Ctrl+C to cancel\n")
        sys.stdout.flush()  # Ensure initial messages are displayed
        
        last_update = time.time()
        update_interval = 0.1  # Update display every 100ms
        current_files = 0
        start_time = time.time()
        
        # Scan each starting path
        for start_path in start_paths:
            for root, dirs, files in os.walk(start_path, onerror=lambda e: stats.add_inaccessible_location(str(e), "Walk error")):
                try:
                    # Remove excluded directories from dirs list to prevent walking them
                    dirs[:] = [d for d in dirs if not is_excluded_path(os.path.join(root, d))]
                    
                    # Update directory count
                    stats.dirs_scanned += 1
                    
                    # Show current progress
                    current_time = time.time()
                    if current_time - last_update >= update_interval:
                        elapsed = current_time - start_time
                        files_per_sec = current_files / elapsed if elapsed > 0 else 0
                        
                        # Clear previous lines and show progress
                        print("\033[K", end="\r")  # Clear line
                        print(f"Scanning: {root}")
                        print(f"Files found: {stats.total_files:,}")
                        print(f"Files indexed: {current_files:,} ({files_per_sec:.0f} files/s)")
                        print("\033[3A", end="")  # Move cursor up 3 lines
                        sys.stdout.flush()  # Force update the display
                        last_update = current_time
                    
                    # Process files in current directory
                    for file in files:
                        file_path = os.path.join(root, file)
                        try:
                            if not is_excluded_path(file_path):
                                # Update current file in stats
                                stats.current_file = file_path
                                stats.total_files += 1
                                current_files += 1
                                
                                try:
                                    # Try to get file size
                                    file_size = os.path.getsize(file_path)
                                    stats.total_bytes += file_size
                                    
                                    # Try to scan the file
                                    safe_scan_file(file_path)
                                    stats.update(bytes_read=file_size, file_completed=True)
                                    sys.stdout.flush()
                                    
                                except OSError as e:
                                    stats.add_inaccessible_location(file_path, f"Cannot access file: {str(e)}")
                                    continue
                                    
                        except Exception as e:
                            stats.add_inaccessible_location(file_path, f"Error processing file: {str(e)}")
                            continue
                            
                except Exception as e:
                    stats.add_inaccessible_location(root, f"Error accessing directory: {str(e)}")
                    continue
        
        # Clear progress display
        print("\033[K", end="\r")  # Clear line
        print("\033[K", end="\r")  # Clear line
        print("\033[K", end="\r")  # Clear line
        
        # Show summary of indexed files
        elapsed = time.time() - start_time
        print(f"\nIndexing Complete!")
        print(f"Time taken: {timedelta(seconds=int(elapsed))}")
        print(f"Total directories found: {stats.dirs_scanned:,}")
        print(f"Total files found: {stats.total_files:,}")
        print(f"Total size: {stats.total_bytes/1024/1024:.1f}MB")
        print(f"Average indexing speed: {stats.total_files/elapsed:.0f} files/s")
        
        if stats.inaccessible_locations:
            print(f"\nWarning: {len(stats.inaccessible_locations):,} locations were inaccessible")
            print("Check inaccessible_locations.txt for details")
        
        print("\nStarting file scanning...")
        sys.stdout.flush()
        
    except KeyboardInterrupt:
        print("\n\nIndexing cancelled by user")
        sys.exit(1)
    finally:
        # Display final statistics
        stats.finalize()

def scan_database(file_path: str) -> None:
    """
    Scan a SQLite database file for potential seed phrases.
    Handles browser cookies and other SQLite databases.
    """
    try:
        import sqlite3
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        
        # Get all table names
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = cursor.fetchall()
        
        for table in tables:
            table_name = table[0]
            try:
                # Get all text columns from the table
                cursor.execute(f"PRAGMA table_info({table_name});")
                columns = cursor.fetchall()
                text_columns = [col[1] for col in columns if 'text' in str(col[2]).lower() or 'varchar' in str(col[2]).lower() or 'char' in str(col[2]).lower()]
                
                # Scan each text column
                for column in text_columns:
                    try:
                        cursor.execute(f"SELECT {column} FROM {table_name} WHERE {column} IS NOT NULL;")
                        rows = cursor.fetchall()
                        for row in rows:
                            if row[0] and isinstance(row[0], str):
                                scan_text(row[0], f"{file_path}:{table_name}.{column}")
                    except sqlite3.Error:
                        continue
            except sqlite3.Error:
                continue
                
        conn.close()
    except (sqlite3.Error, Exception) as e:
        log_unreadable_file(file_path, f"Database error: {str(e)}")
        stats.unreadable_files += 1

def print_config_info():
    """Display current configuration settings"""
    print("\nCurrent Configuration")
    print("====================")
    print("\nScan Settings:")
    print(f"- Operating System: {OPERATING_SYSTEM}")
    print(f"- Min File Size: {config['min_file_size']/1024:.1f}KB")
    print(f"- Max File Size: {config['max_file_size']/1024/1024:.1f}MB")
    
    print("\nTime Filter:")
    if config['use_mtime_filter']:
        filter_time = datetime.fromtimestamp(config['min_mtime'])
        print(f"- Enabled: Only scanning files modified after {filter_time.strftime('%Y-%m-%d %H:%M:%S')}")
    else:
        print("- Disabled: Scanning all files regardless of modification time")
    
    print("\nExcluded Directories:")
    if config['excluded_directories']:
        for directory in config['excluded_directories']:
            print(f"- {directory}")
    else:
        print("- None (scanning all directories)")
    
    print("\nHigh Probability Locations:")
    for directory in config['high_probability_directories']:
        print(f"- {directory}")
    
    print("\nPDF Settings:")
    print(f"- Skip PDF Books: {config['exclude_pdf_books']}")
    if config['exclude_pdf_books']:
        print(f"- PDF Book Threshold: {config['pdf_book_pages']} pages")
    
    print("\nPress Enter to continue...")
    input()

class DatabaseManager:
    """Manages the SQLite database for tracking processed files"""
    def __init__(self):
        self.db_path = 'scan_history.db'
        self._init_db()
    
    def _init_db(self):
        """Initialize the database and create required tables"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            # Create table for processed files
            cursor.execute('''
                CREATE TABLE IF NOT EXISTS processed_files (
                    file_path TEXT PRIMARY KEY,
                    file_id TEXT,
                    scan_time REAL,
                    last_modified REAL
                )
            ''')
            
            # Create table for scan progress
            cursor.execute('''
                CREATE TABLE IF NOT EXISTS scan_progress (
                    id INTEGER PRIMARY KEY,
                    start_time REAL,
                    total_files INTEGER,
                    processed_files INTEGER,
                    completed BOOLEAN
                )
            ''')
            
            conn.commit()
            conn.close()
        except Exception as e:
            logger.error(f"Error initializing database: {str(e)}")
    
    def was_processed(self, file_path: str, file_id: str) -> bool:
        """Check if a file was already processed"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cursor.execute(
                'SELECT file_id FROM processed_files WHERE file_path = ?',
                (file_path,)
            )
            result = cursor.fetchone()
            
            conn.close()
            
            if result is None:
                return False
            
            # If file_id changed, file was modified
            return result[0] == file_id
            
        except Exception as e:
            logger.error(f"Error checking processed status: {str(e)}")
            return False
    
    def add_processed_file(self, file_path: str, file_id: str, last_modified: float):
        """Add a processed file to the database"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cursor.execute('''
                INSERT OR REPLACE INTO processed_files 
                (file_path, file_id, scan_time, last_modified)
                VALUES (?, ?, ?, ?)
            ''', (file_path, file_id, time.time(), last_modified))
            
            conn.commit()
            conn.close()
        except Exception as e:
            logger.error(f"Error adding processed file: {str(e)}")
    
    def get_last_scan_progress(self) -> dict:
        """Get the progress of the last incomplete scan"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cursor.execute('''
                SELECT start_time, total_files, processed_files
                FROM scan_progress
                WHERE completed = 0
                ORDER BY start_time DESC
                LIMIT 1
            ''')
            
            result = cursor.fetchone()
            conn.close()
            
            if result:
                return {
                    'start_time': result[0],
                    'total_files': result[1],
                    'processed_files': result[2]
                }
            return None
            
        except Exception as e:
            logger.error(f"Error getting scan progress: {str(e)}")
            return None
    
    def cleanup_old_records(self, days: int = 30):
        """Remove records older than specified days"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cutoff_time = time.time() - (days * 24 * 3600)
            
            cursor.execute(
                'DELETE FROM processed_files WHERE scan_time < ?',
                (cutoff_time,)
            )
            
            conn.commit()
            conn.close()
        except Exception as e:
            logger.error(f"Error cleaning up old records: {str(e)}")

# Helper functions for file handling
def is_compressed_file(file_path: str) -> bool:
    """Check if a file is a compressed archive"""
    compressed_extensions = {'.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'}
    return any(file_path.lower().endswith(ext) for ext in compressed_extensions)

def scan_zip(file_path: str) -> None:
    """Scan contents of a ZIP file"""
    try:
        with zipfile.ZipFile(file_path) as zf:
            for info in zf.infolist():
                if not info.filename.endswith('/'):  # Skip directories
                    ext = os.path.splitext(info.filename)[1].lower()
                    if ext in target_extensions:
                        try:
                            with zf.open(info) as f:
                                content = f.read().decode('utf-8', errors='ignore')
                                scan_text(content, f"{file_path}:{info.filename}")
                        except Exception as e:
                            logger.debug(f"Error reading {info.filename} in {file_path}: {str(e)}")
    except Exception as e:
        logger.error(f"Error scanning ZIP file {file_path}: {str(e)}")

def is_high_probability_directory(path: str) -> bool:
    """Check if a directory is likely to contain sensitive information"""
    path_lower = path.lower()
    return any(name.lower() in path_lower for name in config['high_probability_directories'])

def log_high_probability_location(path: str) -> None:
    """Log a high probability location for later review"""
    try:
        with open(stats.high_prob_locations_file, 'a', encoding='utf-8') as f:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            f.write(f"{timestamp} | {path}\n")
    except Exception as e:
        logger.error(f"Error logging high probability location: {str(e)}")

def print_startup_info():
    """Display startup banner and information"""
    print("\nBIP39 Seed Phrase Scanner")
    print("========================")
    print("This tool scans files for potential BIP39 seed phrases.")
    print("It can process various file types including:")
    print("- Text files (txt, log, config)")
    print("- Documents (PDF, DOC, DOCX, ODT)")
    print("- Spreadsheets (XLS, XLSX, ODS)")
    print("- Presentations (PPT, PPTX, ODP)")
    print("- Archives (ZIP, RAR, 7Z)")
    print("- Databases (SQLite)")
    print("\nPress Ctrl+C at any time to stop the scan.")

if __name__ == "__main__":
    try:
        # Show startup information
        print_startup_info()
        
        # Load wordlist once at startup
        print("\nLoading wordlist...")
        WORDLIST = load_wordlist()
        BIP39_WORDS = WORDLIST  # Set BIP39_WORDS for compatibility
        print("Wordlist loaded successfully!")
        
        # Validate environment before starting
        if not validate_environment():
            print("\nEnvironment validation failed. Please fix the above issues and try again.")
            sys.exit(1)
        
        # Get OS selection
        OPERATING_SYSTEM = get_os_selection()
        
        # Get excluded directories confirmation
        config['excluded_directories'] = confirm_excluded_directories(OPERATING_SYSTEM)
        
        # Get time filter configuration
        get_time_filter()
        
        # Initialize logging
        setup_logger()
        
        # Print configuration and start scan
        print_config_info()
        main()
    except Exception as e:
        logger.error(f"Fatal error: {str(e)}")
        print(f"\nFatal error occurred: {str(e)}")
        print("Check the log file for details.")
        sys.exit(1)
